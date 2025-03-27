import hashlib
import os
import traceback
from typing import List
import marqo
import pandas as pd
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation

from common import BASE_NAME


def chunk_text(text, chunk_size=300, overlap=50):
    """
    Splits a text into chunks of chunk_size words with overlap.

    :param text: The full document text
    :param chunk_size: Number of words per chunk
    :param overlap: Number of overlapping words between chunks
    :return: List of text chunks
    """
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i : i + chunk_size])
        chunks.append(chunk)
    return chunks


# Read TXT
def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return chunk_text(f.read())


# Read DOCX
def read_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return chunk_text(text)


# Read PDF
def read_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        text = "\n".join(
            page.extract_text() for page in pdf.pages if page.extract_text()
        )
        return chunk_text(text)


# Read PPTX
def read_pptx(file_path):
    ppt = Presentation(file_path)
    text = "\n".join(
        [
            shape.text
            for slide in ppt.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ]
    )
    return chunk_text(text)


# Read Excel
def read_excel(file_path):
    df = pd.read_excel(file_path, engine="openpyxl")
    return [str(row) for row in df.to_dict(orient="records")]


# Extract Text from Images (OCR)
def read_image(file_path):
    image = Image.open(file_path)
    text = "".join(pytesseract.image_to_string(image))
    print(text)
    return chunk_text(text)


def index_folder(folder_path):

    for dirpath, _, filenames in os.walk(folder_path):
        for file in filenames:
            documents = []
            file_path = os.path.join(dirpath, file)
            file_ext = file.split(".")[-1].lower()
            chunks = []
            print(f"Indexing {file_path}")

            try:
                if file_ext == "txt":
                    chunks = read_txt(file_path)
                elif file_ext == "docx":
                    chunks = read_docx(file_path)
                elif file_ext == "pdf":
                    chunks = read_pdf(file_path)
                elif file_ext == "pptx":
                    chunks = read_pptx(file_path)
                elif file_ext == "xlsx":
                    chunks = read_excel(file_path)
                elif file_ext in ["jpg", "png", "jpeg"]:
                    chunks = read_image(file_path)  # OCR extraction
            except Exception as ex:
                print_error(ex)
                traceback.print_exc()
                continue

            # Add each chunk separately
            for i, chunk in enumerate(chunks):
                hash_id = hashlib.md5(f"{file_path}_chunk_{i}".encode()).hexdigest()
                item = {
                    "_id": hash_id,
                    "title": file,
                    "content": chunk,
                    "chunk_name": f"{file}_chunk_{i}",
                    "file_type": file_ext,
                    "file_path": file_path,
                }
                documents.append(item)

            # # Delete old chunk based on file_path if needed
            # existing_chunks = mq.index(BASE_NAME).search(
            #     file_path, filter_string=f'file_path:"{file_path}"'
            # )
            # if existing_chunks["hits"]:
            #     ids_to_delete = [hit["_id"] for hit in existing_chunks["hits"]]
            #     mq.index(BASE_NAME).delete_documents(ids_to_delete)
            #     print(f"Deleted {len(ids_to_delete)} old chunks of {file_path}")

            # Index chunks
            if documents:
                batch_index_documents(documents)


def batch_index_documents(documents, batch_size=128):
    """Splits documents into batches of 128 and indexes them in Marqo"""
    for i in range(0, len(documents), batch_size):
        batch = documents[i : i + batch_size]  # Get a batch of 128 or less
        result = mq.index(BASE_NAME).add_documents(
            batch, tensor_fields=["content", "title"]
        )
        if result["errors"]:
            print_error(result)
            print(documents)
            raise Exception("Add document failed")
        print(f"Indexed batch {i//batch_size + 1}/{(len(documents) // batch_size) + 1}")


def print_error(content):
    print(f"\033[31m{content}\033[0m")


mq = marqo.Client()

folder_path = "/Users/thangnguyen/Downloads/DooDoo Files"
index_folder(folder_path)

results = mq.index(BASE_NAME).search(q="What is DooDoo?")
print(results)
